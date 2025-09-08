# AI CSV Interpreter v3 — extended (adds built-in stats, regression, DOE suggestions, PDF export, xlsx support)
# Keep original features: Gemini integration, safe python execution, slide editor, live preview, PPTX export.

import os
os.environ["PATH"] += os.pathsep + r"C:\Program Files\LibreOffice\program"
os.environ["PATH"] += os.pathsep + r"C:\Program Files\poppler-24.08.0\Library\bin"

import io, contextlib, subprocess, tempfile, time, shutil, traceback, re, json, base64
from dotenv import load_dotenv

import streamlit as st
import streamlit.components.v1 as components
import pandas as pd
import plotly.express as px
import google.generativeai as genai
import matplotlib.pyplot as plt
import seaborn as sns
from pptx import Presentation
from pptx.util import Inches, Pt
from PIL import Image, ImageDraw, ImageFont
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN

# Optional libs
try:
    from pdf2image import convert_from_path
    PDF2IMAGE_AVAILABLE = True
except Exception:
    PDF2IMAGE_AVAILABLE = False

try:
    from streamlit_autorefresh import st_autorefresh
    AUTOREFRESH_AVAILABLE = True
except Exception:
    AUTOREFRESH_AVAILABLE = False

# Stats & ML optional libs
try:
    import scipy.stats as stats
except Exception:
    stats = None

try:
    import statsmodels.api as sm
    from statsmodels.formula.api import ols
    STATSMODELS_AVAILABLE = True
except Exception:
    STATSMODELS_AVAILABLE = False

try:
    from sklearn.linear_model import LinearRegression
    from sklearn.preprocessing import StandardScaler
    SKLEARN_AVAILABLE = True
except Exception:
    SKLEARN_AVAILABLE = False

try:
    from reportlab.lib.pagesizes import letter
    from reportlab.pdfgen import canvas
    REPORTLAB_AVAILABLE = True
except Exception:
    REPORTLAB_AVAILABLE = False

# ---------------------- utils ----------------------
def _which(cmd):
    return shutil.which(cmd) is not None

SOFFICE_OK = _which("soffice")
POPPLER_OK = _which("pdftoppm")

# -------------------- session state --------------------
st.set_page_config(page_title="AI CSV Interpreter v3 — Slide Editor", layout="wide")
ss = st.session_state

# core session keys
if "ppt" not in ss: ss.ppt = Presentation()
if "messages" not in ss: ss.messages = []          # list of dicts: {role, content}
if "df" not in ss: ss.df = None
if "msg_plot_map" not in ss: ss.msg_plot_map = {}   # {msg_idx: img_path}
if "last_plot_path" not in ss: ss.last_plot_path = None
if "preview_images" not in ss: ss.preview_images = []
if "preview_dirty" not in ss: ss.preview_dirty = True
if "autorefresh_on" not in ss: ss.autorefresh_on = False
if "refresh_interval_ms" not in ss: ss.refresh_interval_ms = 5000
if "slide_idx" not in ss: ss.slide_idx = 1

# Slide Editor state
if "slide_editor_title" not in ss: ss.slide_editor_title = "Slide Title"
if "slide_editor_text" not in ss: ss.slide_editor_text = ""
if "slide_editor_selected_plot" not in ss: ss.slide_editor_selected_plot = None

# -------------------- New Slide Layout Controls --------------------
if "layout_style" not in ss: ss.layout_style = "Text + Image (side-by-side)"
if "font_size_pt" not in ss: ss.font_size_pt = 14
if "img_width_in" not in ss: ss.img_width_in = 5.0
if "img_height_in" not in ss: ss.img_height_in = 3.0
if "reuse_layout" not in ss: ss.reuse_layout = False

# -------------------- Gemini config --------------------
load_dotenv()
GOOGLE_API_KEY = os.getenv("GOOGLE_API_KEY")
st.title("📊 AI CSV Interpreter v3 — Slide Editor + PPT Export")
st.caption(f"LibreOffice: {'✅' if SOFFICE_OK else '❌'} | Poppler: {'✅' if POPPLER_OK else '❌'} | pdf2image: {'✅' if PDF2IMAGE_AVAILABLE else '❌'}")
if not GOOGLE_API_KEY:
    st.error("Missing GOOGLE_API_KEY in .env. Add GOOGLE_API_KEY=... and restart.")
    st.stop()
genai.configure(api_key=GOOGLE_API_KEY)
model = genai.GenerativeModel("gemini-1.5-flash-latest")

# -------------------- helpers (existing + new) --------------------
def read_csv_with_fallback(file):
    encodings = ['utf-8', 'latin1', 'ISO-8859-1', 'cp1252']
    for e in encodings:
        try:
            file.seek(0)
            return pd.read_csv(file, encoding=e)
        except Exception:
            continue
    raise Exception("Could not decode CSV with common encodings.")

def read_table_auto(file, name):
    """Support csv and xlsx"""
    if name.lower().endswith(".csv"):
        return read_csv_with_fallback(file)
    elif name.lower().endswith((".xls", ".xlsx")):
        try:
            file.seek(0)
            return pd.read_excel(file)
        except Exception as e:
            raise
    else:
        return read_csv_with_fallback(file)

def _save_matplotlib(msg_idx=None):
    if not plt.get_fignums():
        return None
    save_dir = os.path.join(tempfile.gettempdir(), "ai_saved_plots")
    os.makedirs(save_dir, exist_ok=True)
    path = os.path.join(save_dir, f"mpl_plot_msg{msg_idx or 'X'}_{int(time.time()*1000)}.png")
    plt.savefig(path, dpi=150, bbox_inches='tight')
    plt.close('all')
    return path

def _try_save_plotly(user_ns, msg_idx=None):
    try:
        fig = user_ns.get("fig", None)
        if fig is None:
            return None
        save_dir = os.path.join(tempfile.gettempdir(), "ai_saved_plots")
        os.makedirs(save_dir, exist_ok=True)
        path = os.path.join(save_dir, f"plotly_msg{msg_idx or 'X'}_{int(time.time()*1000)}.png")
        if hasattr(fig, "write_image"):
            try:
                fig.write_image(path, scale=2)
                return path
            except Exception:
                pass
        import plotly.io as pio
        try:
            pio.write_image(fig, path, scale=2)
            return path
        except Exception:
            return None
    except Exception:
        return None

def run_generated_code(code, df, msg_idx=None):
    """
    Execute code in isolated namespace. Capture stdout and matplotlib/plotly figures.
    Returns (stdout_text, error_text).
    """
    out_buf = io.StringIO()
    user_ns = {"df": df, "pd": pd, "plt": plt, "sns": sns, "px": px, "np": __import__("numpy")}
    err = None
    with contextlib.redirect_stdout(out_buf):
        try:
            exec(code, user_ns)
        except Exception:
            err = traceback.format_exc()
    printed = out_buf.getvalue()
    try:
        plot_path = None
        if plt.get_fignums():
            plot_path = _save_matplotlib(msg_idx)
        else:
            plot_path = _try_save_plotly(user_ns, msg_idx)
        if plot_path:
            ss.last_plot_path = plot_path
            if msg_idx is not None:
                ss.msg_plot_map[msg_idx] = plot_path
    except Exception as e:
        err = (err + f"\nAdditionally failed saving plot: {e}") if err else f"Failed saving plot: {e}"
    return printed, err

# ---------- new analysis & export helpers ----------

def perform_descriptive_stats(df):
    out = {}
    try:
        out['describe'] = df.describe(include='all').to_dict()
        numeric = df.select_dtypes(include='number')
        out['ncols'] = numeric.shape[1]
    except Exception as e:
        out['error'] = str(e)
    return out

def pearson_corr_with_pvalues(df):
    """Return DataFrame of correlations and p-values for numeric columns"""
    num = df.select_dtypes(include='number')
    cols = num.columns.tolist()
    n = len(cols)
    corr = pd.DataFrame(index=cols, columns=cols, dtype=float)
    pvals = pd.DataFrame(index=cols, columns=cols, dtype=float)
    for i in range(n):
        for j in range(n):
            x = num.iloc[:, i].dropna()
            y = num.iloc[:, j].dropna()
            # align indices by index intersection
            xy = pd.concat([num.iloc[:, i], num.iloc[:, j]], axis=1).dropna()
            if xy.shape[0] < 3 or (xy.iloc[:,0].std()==0 and xy.iloc[:,1].std()==0):
                corr.iloc[i,j] = float('nan')
                pvals.iloc[i,j] = float('nan')
                continue
            try:
                r, p = stats.pearsonr(xy.iloc[:,0], xy.iloc[:,1]) if stats else (None,None)
                corr.iloc[i,j] = r
                pvals.iloc[i,j] = p
            except Exception:
                corr.iloc[i,j] = float('nan')
                pvals.iloc[i,j] = float('nan')
    return corr, pvals

def auto_select_dependent(df):
    """Pick dependent var heuristically: highest variance or last numeric column"""
    num = df.select_dtypes(include='number')
    if num.shape[1] == 0:
        return None
    try:
        variances = num.var().sort_values(ascending=False)
        return variances.index[0]
    except Exception:
        return num.columns[-1]

def run_regression(df, dep=None):
    """Run linear regression and return summary dict"""
    num = df.select_dtypes(include='number').dropna()
    if dep is None:
        dep = auto_select_dependent(df)
    if dep not in num.columns:
        return {"error":"No numeric dependent variable found."}
    X = num.drop(columns=[dep])
    y = num[dep]
    if X.shape[1] == 0:
        return {"error":"No independent numeric columns found for regression."}
    try:
        if STATSMODELS_AVAILABLE:
            # statsmodels requires adding constant
            X2 = sm.add_constant(X)
            model_sm = sm.OLS(y, X2).fit()
            return {"summary": model_sm.summary().as_text(), "params": model_sm.params.to_dict(), "rsquared": float(model_sm.rsquared)}
        elif SKLEARN_AVAILABLE:
            lr = LinearRegression()
            lr.fit(X.fillna(0), y.fillna(0))
            params = dict(zip(X.columns, lr.coef_.tolist()))
            return {"params": params, "intercept": float(lr.intercept_), "rsquared": None}
        else:
            return {"error":"No regression libs available (install statsmodels or scikit-learn)."}
    except Exception as e:
        return {"error": str(e)}

def run_hypothesis_tests(df):
    """Run candidate tests: t-test between groups if a 'group' column exists; Levene for variance differences; ANOVA if >2 groups"""
    out = {}
    num = df.select_dtypes(include='number')
    # Try to detect categorical grouping column
    cats = [c for c in df.columns if df[c].nunique() < 10 and df[c].dtype == object]
    if not cats:
        # try integer-coded groups
        cats = [c for c in df.columns if df[c].nunique() < 10 and pd.api.types.is_integer_dtype(df[c].dtype)]
    if cats:
        gcol = cats[0]
        groups = df.groupby(gcol)
        if len(groups) >= 2:
            # choose first numeric column for example
            if num.shape[1] > 0:
                col = num.columns[0]
                samples = [grp[col].dropna().values for _, grp in groups]
                try:
                    if len(samples) == 2:
                        tstat, p = stats.ttest_ind(samples[0], samples[1], equal_var=False)
                        out['t_test'] = {"column":col, "tstat": float(tstat), "pvalue": float(p), "groups": list(groups.groups.keys())}
                    elif len(samples) > 2:
                        f, p = stats.f_oneway(*samples)
                        out['anova'] = {"column":col, "fstat": float(f), "pvalue": float(p)}
                    # Levene
                    w, p_levene = stats.levene(*samples)
                    out['levene'] = {"w": float(w), "pvalue": float(p_levene)}
                except Exception as e:
                    out['error'] = str(e)
    else:
        out['note'] = "No categorical grouping column detected to run t-test/ANOVA."
    return out

def suggest_next_experiments(df, target_col=None, top_k=3):
    """Heuristic suggestions: select top correlated inputs and suggest +/- perturbations."""
    suggestions = []
    num = df.select_dtypes(include='number')
    if num.shape[1] < 2:
        return {"error": "Not enough numeric columns to make suggestions."}
    if target_col is None:
        target_col = auto_select_dependent(df)
    if target_col not in num.columns:
        return {"error": "Target column not found."}
    # correlation
    try:
        corrs = num.corr()[target_col].abs().sort_values(ascending=False)
        corrs = corrs.drop(index=[target_col])
        tops = corrs.head(top_k).index.tolist()
        for col in tops:
            mean = num[col].mean()
            std = num[col].std()
            # suggest testing mean +/- 0.5*std
            suggestions.append({
                "variable": col,
                "current_mean": float(mean),
                "suggest_test_values": [float(max(mean - 0.5*std, mean*0.5)), float(mean + 0.5*std)]
            })
        return {"target": target_col, "suggestions": suggestions}
    except Exception as e:
        return {"error": str(e)}

def _img_to_data_uri(path: str) -> str:
    try:
        with open(path, "rb") as f:
            b = f.read()
        return f"data:image/png;base64,{base64.b64encode(b).decode('ascii')}"
    except Exception:
        return ""

# PDF export (reportlab or pillow fallback)
def generate_pdf_report(text_blocks, image_paths, out_path=None, title="AI Report"):
    if out_path is None:
        out_path = tempfile.NamedTemporaryFile(delete=False, suffix=".pdf").name
    try:
        if REPORTLAB_AVAILABLE:
            c = canvas.Canvas(out_path, pagesize=letter)
            width, height = letter
            # Title
            c.setFont("Helvetica-Bold", 16)
            c.drawString(36, height - 50, title)
            y = height - 80
            c.setFont("Helvetica", 10)
            for block in text_blocks:
                for line in block.splitlines():
                    if y < 80:
                        c.showPage()
                        y = height - 50
                    c.drawString(36, y, line[:120])
                    y -= 12
                y -= 8
            # images appended on new pages
            for imgp in image_paths:
                try:
                    c.showPage()
                    c.drawImage(imgp, 36, 72, width=width-72, preserveAspectRatio=True, anchor='c')
                except Exception:
                    pass
            c.save()
            return out_path
        else:
            # Pillow fallback: make multipage PDF by converting images and text to images
            pages = []
            W, H = 1100, 1400
            fnt = ImageFont.load_default()
            # first page: text
            img = Image.new("RGB", (W, H), color=(255,255,255))
            d = ImageDraw.Draw(img)
            y = 40
            d.text((36, y), title, font=fnt, fill=(0,0,0))
            y += 30
            for block in text_blocks:
                for line in block.splitlines():
                    d.text((36, y), line[:120], font=fnt, fill=(0,0,0))
                    y += 14
                    if y > H - 100:
                        pages.append(img)
                        img = Image.new("RGB", (W, H), color=(255,255,255))
                        d = ImageDraw.Draw(img)
                        y = 40
            pages.append(img)
            # image pages
            for pth in image_paths:
                try:
                    im = Image.open(pth).convert("RGB")
                    im.thumbnail((W-72, H-200))
                    bg = Image.new("RGB", (W, H), color=(255,255,255))
                    bg.paste(im, (36,100))
                    pages.append(bg)
                except Exception:
                    pass
            pages[0].save(out_path, save_all=True, append_images=pages[1:])
            return out_path
    except Exception as e:
        raise

def _html_viewer_template(img_data_uri: str, slide_idx: int, total: int) -> str:
    """Small, safe HTML template that displays an image and slide counter."""
    return f"""
    <div style="display:flex;flex-direction:column;align-items:center;justify-content:center;font-family:Arial,Helvetica,sans-serif;">
      <div style="width:100%;display:flex;justify-content:center;">
        <img src="{img_data_uri}" style="max-width:100%; height:auto; border-radius:8px; box-shadow:0 6px 18px rgba(0,0,0,0.12);" />
      </div>
      <div style="margin-top:8px; color:#333; font-size:14px;">
        Slide <strong>{slide_idx}</strong> of <strong>{total}</strong>
      </div>
    </div>
    """

def _render_floating_viewer(preview_images, slide_idx=1, height=420):
    """
    Render a single preview image (as an embedded data URI) using Streamlit components.html.
    preview_images: list of filesystem paths to PNG images
    slide_idx: 1-based index of which image to render
    height: iframe height in px
    """
    try:
        if not preview_images:
            # nothing to show
            components.html("<div style='color:#666;padding:12px;'>No preview images available.</div>", height=120)
            return

        # clamp slide_idx into valid range
        total = len(preview_images)
        idx = min(max(1, int(slide_idx)), total) - 1
        img_path = preview_images[idx] if idx < total else preview_images[0]

        # convert to data URI (uses the existing helper _img_to_data_uri)
        img_data_uri = _img_to_data_uri(img_path)
        if not img_data_uri:
            components.html("<div style='color:#a00;padding:12px;'>Preview image missing or unreadable.</div>", height=120)
            return

        html = _html_viewer_template(img_data_uri, idx + 1, total)
        components.html(html, height=height)

    except Exception as e:
        # fail-safe: show error text in UI
        components.html(f"<div style='color:red;padding:8px;'>Preview render error: {str(e)[:300]}</div>", height=120)

def _render_floating_viewer(preview_images, slide_idx=1, height=420):
    """
    Render a preview image using Streamlit components.
    Handles both cases where _html_viewer_template is a function or a string template.
    """
    try:
        if not preview_images:
            components.html("<div style='color:#666;padding:12px;'>No preview images available.</div>", height=120)
            return

        total = len(preview_images)
        idx = min(max(1, int(slide_idx)), total) - 1
        img_path = preview_images[idx] if idx < total else preview_images[0]

        # convert to data URI for embedding
        img_data_uri = _img_to_data_uri(img_path)
        if not img_data_uri:
            components.html("<div style='color:#a00;padding:12px;'>Preview image missing or unreadable.</div>", height=120)
            return

        # Build HTML safely: call template if it's a function, otherwise treat as string and replace placeholders
        if callable(_html_viewer_template):
            html = _html_viewer_template(img_data_uri, idx + 1, total)
        else:
            # legacy string template expected to contain tokens: __DATA__, __START__, __MAX__, __VAL__
            html = str(_html_viewer_template)
            html = html.replace("__DATA__", img_data_uri) \
                       .replace("__START__", str(idx)) \
                       .replace("__MAX__", str(total - 1)) \
                       .replace("__VAL__", str(idx))

        components.html(html, height=height)

    except Exception as e:
        components.html(f"<div style='color:red;padding:8px;'>Preview render error: {str(e)[:300]}</div>", height=120)
        return

# -------------------- Updated add_slide_with_text_and_optional_image (kept identical logic) --------------------
def add_slide_with_text_and_optional_image(title, text, image_path=None):
    layout = ss.ppt.slide_layouts[1] if len(ss.ppt.slide_layouts) > 1 else ss.ppt.slide_layouts[5]
    slide = ss.ppt.slides.add_slide(layout)

    # Title
    try:
        slide.shapes.title.text = title
    except Exception:
        slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.5)).text_frame.text = title

    # Handle layouts
    style = ss.layout_style
    font_size = Pt(ss.font_size_pt)

    if style == "Text only":
        tx = slide.shapes.add_textbox(Inches(0.5), Inches(1), Inches(9), Inches(5))
        tf = tx.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = font_size

    elif style == "Image only" and image_path and os.path.exists(image_path):
        slide.shapes.add_picture(image_path, Inches(1), Inches(1), width=Inches(ss.img_width_in), height=Inches(ss.img_height_in))

    elif style == "Text + Image (side-by-side)":
        # Text left
        tx = slide.shapes.add_textbox(Inches(0.5), Inches(1), Inches(4.5), Inches(5))
        tf = tx.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = font_size
        # Image right
        if image_path and os.path.exists(image_path):
            slide.shapes.add_picture(image_path, Inches(5.2), Inches(1), width=Inches(ss.img_width_in), height=Inches(ss.img_height_in))

    elif style == "Text top + Image bottom":
        # Text on top
        tx = slide.shapes.add_textbox(Inches(0.5), Inches(1), Inches(9), Inches(2))
        tf = tx.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = font_size
        # Image bottom
        if image_path and os.path.exists(image_path):
            slide.shapes.add_picture(image_path, Inches(1), Inches(3.2), width=Inches(ss.img_width_in), height=Inches(ss.img_height_in))

    elif style == "2x2 Image Grid" and image_path and os.path.exists(image_path):
        # duplicate same image into 4 slots
        w, h = ss.img_width_in/2, ss.img_height_in/2
        slide.shapes.add_picture(image_path, Inches(0.5), Inches(1), width=Inches(w), height=Inches(h))
        slide.shapes.add_picture(image_path, Inches(5), Inches(1), width=Inches(w), height=Inches(h))
        slide.shapes.add_picture(image_path, Inches(0.5), Inches(3.5), width=Inches(w), height=Inches(h))
        slide.shapes.add_picture(image_path, Inches(5), Inches(3.5), width=Inches(w), height=Inches(h))

    ss.preview_dirty = True

def generate_live_preview_images():
    """
    Render the current ss.ppt slides to image files and return list of image paths.
    Tries:
      - Save PPTX -> convert to PDF with soffice -> convert PDF pages to images via pdf2image
      - Fallback: generate simple PNG preview images using slide text (PIL)
    """
    imgs = []
    try:
        # save pptx to temp file
        tmp_ppt = tempfile.NamedTemporaryFile(delete=False, suffix=".pptx")
        ss.ppt.save(tmp_ppt.name)
        tmp_ppt.close()

        # Try LibreOffice conversion PPTX -> PDF
        pdf_path = tmp_ppt.name.replace(".pptx", ".pdf")
        if SOFFICE_OK:
            try:
                # Convert to PDF (headless)
                cmd = ["soffice", "--headless", "--convert-to", "pdf", "--outdir", os.path.dirname(pdf_path), tmp_ppt.name]
                subprocess.run(cmd, check=True, stdout=subprocess.PIPE, stderr=subprocess.PIPE, timeout=30)
                # If pdf2image available, convert PDF pages -> images
                if PDF2IMAGE_AVAILABLE and os.path.exists(pdf_path):
                    pages = convert_from_path(pdf_path, dpi=150)
                    for i, page in enumerate(pages):
                        out_png = tempfile.NamedTemporaryFile(delete=False, suffix=f"_slide{i+1}.png").name
                        page.save(out_png, "PNG")
                        imgs.append(out_png)
                    # cleanup temp files if desired
                    return imgs
            except Exception:
                # fallback path if conversion failed
                pass

        # If conversion not possible, attempt to create thumbnails from slide data (PIL fallback)
        from PIL import Image, ImageDraw, ImageFont

        def _render_slide_to_image(slide, idx):
            W, H = 1200, 900
            bg = Image.new("RGB", (W, H), color=(255, 255, 255))
            draw = ImageDraw.Draw(bg)
            # draw title if present
            title_text = ""
            body_text = ""
            try:
                if slide.shapes.title:
                    title_text = slide.shapes.title.text if slide.shapes.title.text else ""
            except Exception:
                title_text = ""
            # aggregate first few textboxes
            try:
                texts = []
                for shp in slide.shapes:
                    if shp.shape_type == MSO_SHAPE_TYPE.TEXT_BOX or shp.has_text_frame:
                        try:
                            t = shp.text_frame.text.strip()
                            if t:
                                texts.append(t)
                        except Exception:
                            continue
                body_text = "\n\n".join(texts[:6])
            except Exception:
                body_text = ""
            # fonts (use default if PIL font missing)
            try:
                f_title = ImageFont.truetype("arial.ttf", 28)
                f_body = ImageFont.truetype("arial.ttf", 14)
            except Exception:
                f_title = ImageFont.load_default()
                f_body = ImageFont.load_default()
            # draw title
            if title_text:
                draw.text((36, 24), title_text[:200], font=f_title, fill=(0, 0, 0))
            # draw body
            y = 80
            for line in (body_text or "").splitlines():
                if y > H - 60:
                    break
                draw.text((36, y), line[:200], font=f_body, fill=(0, 0, 0))
                y += 18
            # if slide has images/plots referenced in ss.msg_plot_map, try to paste one
            try:
                # search for an image path that references this slide index (best-effort)
                for p in ss.msg_plot_map.values():
                    if p and os.path.exists(p):
                        thumb = Image.open(p)
                        thumb.thumbnail((int(W*0.6), int(H*0.45)))
                        bg.paste(thumb, (36, max(140, (H - thumb.size[1])//2)))
                        break
            except Exception:
                pass
            out = tempfile.NamedTemporaryFile(delete=False, suffix=f"_fallback_slide{idx+1}.png").name
            bg.save(out, "PNG")
            return out

        # Create images for each slide
        for i, slide in enumerate(ss.ppt.slides):
            try:
                imgpath = _render_slide_to_image(slide, i)
                imgs.append(imgpath)
            except Exception:
                continue

        # if nothing created, create one blank placeholder
        if not imgs:
            W, H = 1200, 900
            img = Image.new("RGB", (W, H), color=(240, 240, 240))
            draw = ImageDraw.Draw(img)
            try:
                f = ImageFont.truetype("arial.ttf", 20)
            except Exception:
                f = ImageFont.load_default()
            draw.text((40, 40), "No slides available", font=f, fill=(80, 80, 80))
            out = tempfile.NamedTemporaryFile(delete=False, suffix="_placeholder.png").name
            img.save(out, "PNG")
            imgs.append(out)

        return imgs

    except Exception as e:
        # Last-resort fallback single image explaining error
        try:
            W, H = 1200, 900
            img = Image.new("RGB", (W, H), color=(255, 255, 255))
            draw = ImageDraw.Draw(img)
            try:
                f = ImageFont.truetype("arial.ttf", 18)
            except Exception:
                f = ImageFont.load_default()
            draw.text((36, 40), "Preview generation failed: " + str(e)[:200], font=f, fill=(255, 0, 0))
            out = tempfile.NamedTemporaryFile(delete=False, suffix="_error.png").name
            img.save(out, "PNG")
            return [out]
        except Exception:
            return []


# (rest of helper functions for PPT preview remain unchanged; omitted for brevity in this message but retained in file)

# -------------------- UI and integration (left/right layout) --------------------
EMU_PER_INCH = 914400
def _emu_to_in(v):
    try: return float(v) / EMU_PER_INCH
    except Exception: return 0.0
def _in_to_emu(v):
    try: return int(float(v) * EMU_PER_INCH)
    except Exception: return 0

# the code that manages the slide editor / ppt editing / preview remains unchanged (kept exactly as in your original)

# -------------------- Main UI panels (left/right) --------------------
left, right = st.columns([2, 1])

with right:
    # Slide Layout Controls (unchanged UI from your script)
    with st.expander("⚙️ Slide Layout Controls", expanded=True):
        ss.layout_style = st.radio(
            "Choose layout",
            ["Text only", "Image only", "Text + Image (side-by-side)", "Text top + Image bottom", "2x2 Image Grid"],
            index=["Text + Image (side-by-side)", "Text only", "Image only", "Text top + Image bottom", "2x2 Image Grid"].index(ss.layout_style)
        )
        ss.font_size_pt = st.slider("Font size (pt)", 8, 32, ss.font_size_pt, step=1)
        ss.img_width_in = st.slider("Image width (inches)", 2.0, 8.0, ss.img_width_in, step=0.5)
        ss.img_height_in = st.slider("Image height (inches)", 2.0, 6.0, ss.img_height_in, step=0.5)
        ss.reuse_layout = st.checkbox("Reuse last layout", value=ss.reuse_layout)

with left:
    # CSV / XLSX upload
    uploaded_file = st.file_uploader("Upload CSV or Excel file", type=["csv","xls","xlsx"], key="uploader_csv")
    if uploaded_file:
        try:
            st.session_state.df = read_table_auto(uploaded_file, uploaded_file.name)
            st.subheader("Preview (first 5 rows)")
            st.dataframe(st.session_state.df.head(5))

            # quick line chart if numeric columns
            numeric_cols = st.session_state.df.select_dtypes(include="number").columns.tolist()
            if len(numeric_cols) >= 2:
                x_axis = st.selectbox("X-axis", numeric_cols, key="xaxis_select")
                y_axis = st.multiselect("Y-Axis", numeric_cols, default=numeric_cols[1:2], key="yaxis_select")
                if x_axis and y_axis:
                    fig = px.line(st.session_state.df, x=x_axis, y=y_axis, title="Line Chart")
                    st.plotly_chart(fig, use_container_width=True)

            # Initial AI analysis (existing Gemini call)
            if st.button("🔍 Run Initial AI Analysis (Gemini)", key="run_initial_ai"):
                prompt = f"""
                You are an expert engineer with great data analysis skills. Depending on data type, if applicable:
                   1. perform a full statistical analysis.
                      I need a report that includes the calculated results of statistical tests,
                      not just a general explanation of how to perform them.
                      Specifically, calculate the mean and variance for each group,
                      then perform a t-test for the difference in means and a Levene's test for the difference in variances.
                      Based on your findings, provide clear, data-driven recommendations for process improvement.

                    2. 🔁 Pearson correlation matrix with p-values. Include statistical interpretation.
                    3. 📈 Multiple linear regression:
                       - Select the most likely dependent variable automatically (or infer).
                       - Show coefficients, R², p-values, and model summary.
                    4. 🧪 Perform appropriate hypothesis testing (e.g., ANOVA or t-test) where applicable.
                Data:
                {st.session_state.df.to_csv(index=False)}
                """
                try:
                    response = model.generate_content(prompt)
                    ss.messages.append({"role": "assistant", "content": response.text})
                    st.success("Gemini analysis complete; results appended to conversation.")
                except Exception as e:
                    st.error("AI call failed: " + str(e))

            # NEW: Built-in deterministic analysis (safe, reproducible)
            st.markdown("### ⚙️ Built-in deterministic analysis")
            colA, colB, colC = st.columns(3)
            with colA:
                if st.button("Run built-in stats & correlation"):
                    df = st.session_state.df
                    try:
                        desc = perform_descriptive_stats(df)
                        st.write("Descriptive statistics:")
                        st.json(desc.get('describe', {}))
                        if stats:
                            corr, pvals = pearson_corr_with_pvalues(df)
                            st.write("Pearson correlation matrix:")
                            st.dataframe(corr)
                            st.write("P-values matrix:")
                            st.dataframe(pvals)
                        else:
                            st.warning("scipy not available; install scipy to compute p-values.")
                    except Exception as e:
                        st.error("Error running built-in stats: " + str(e))
            with colB:
                if st.button("Run regression (linear)"):
                    df = st.session_state.df
                    res = run_regression(df)
                    if 'error' in res:
                        st.error(res['error'])
                    else:
                        st.write("Regression results:")
                        if 'summary' in res:
                            st.text(res['summary'])
                        else:
                            st.json({k:v for k,v in res.items() if k!='summary'})
            with colC:
                if st.button("Run hypothesis tests (t/ANOVA/Levene)"):
                    df = st.session_state.df
                    res = run_hypothesis_tests(df)
                    st.write("Hypothesis test results:")
                    st.json(res)

            st.markdown("---")
            # DOE suggestion
            st.subheader("🔬 DOE / Next-experiment suggestions")
            if st.button("Suggest next experiments"):
                try:
                    df = st.session_state.df
                    target = auto_select_dependent(df)
                    sug = suggest_next_experiments(df, target_col=target)
                    st.write("Suggestions (heuristic):")
                    st.json(sug)
                except Exception as e:
                    st.error("DOE suggestion failed: " + str(e))

            st.markdown("---")
            # Save dataset to history
            hist_dir = os.path.join(os.getcwd(), "datasets")
            os.makedirs(hist_dir, exist_ok=True)
            fname = f"{time.strftime('%Y%m%d_%H%M%S')}_{uploaded_file.name}"
            path = os.path.join(hist_dir, fname)
            try:
                st.session_state.df.to_csv(path, index=False)
                st.info(f"Saved uploaded dataset to history: {path}")
            except Exception:
                pass

            st.markdown("---")
            # PDF export using built-in analysis + optional plots
            if st.button("📄 Export PDF report (built-in analysis + last plot)"):
                try:
                    text_blocks = []
                    text_blocks.append("AI CSV Interpreter — Report")
                    text_blocks.append(f"File: {uploaded_file.name}")
                    # include built-in stats textual summary
                    try:
                        desc = perform_descriptive_stats(st.session_state.df)
                        text_blocks.append("Descriptive statistics (summary):")
                        text_blocks.append(str(desc.get('describe', {})))
                    except Exception:
                        pass
                    # include last plot path
                    images = []
                    if ss.last_plot_path and os.path.exists(ss.last_plot_path):
                        images.append(ss.last_plot_path)
                    pdf_path = generate_pdf_report(text_blocks, images)
                    with open(pdf_path, "rb") as f:
                        st.download_button("Download PDF report", f, file_name="ai_report.pdf", mime="application/pdf")
                except Exception as e:
                    st.error("PDF export failed: " + str(e))

        except Exception as e:
            st.error("Error reading file: " + str(e))

    st.markdown("---")
    # Chat / follow-ups (existing code preserved)
    st.subheader("💬 Ask follow-up questions (chat)")

    try:
        follow = st.chat_input("Ask a follow-up question or request a clarification...")
    except Exception:
        follow = st.text_input("Ask a follow-up question or request a clarification...", key="fallback_follow")

    if follow:
        ss.messages.append({"role": "user", "content": follow})
        # Build compact history; include small sample of dataset optionally
        history = "\n".join([f"{m['role'].capitalize()}: {m['content']}" for m in ss.messages])
        if ss.df is not None:
            try:
                history += "\n\nDataset (first 200 rows):\n" + ss.df.head(200).to_csv(index=False)
            except Exception:
                pass
        # Instruction for code-only answers when user asks for scripts
        history += """
IMPORTANT: If the user requests a Python script, respond ONLY with the Python code inside triple backticks tagged with Python. No extra explanation.
"""
        try:
            resp = model.generate_content(history)
            ss.messages.append({"role": "assistant", "content": resp.text})
            ss.preview_dirty = True
            st.rerun()
        except Exception as e:
            st.error("AI call failed: " + str(e))

    st.markdown("---")
    st.subheader("AI conversation & actions")

    # (display messages and run code blocks — existing code preserved from your script)
    for idx, msg in enumerate(ss.messages):
        if msg["role"] == "user":
            st.markdown(f"**You:** {msg['content']}")
        else:
            st.markdown(f"**AI #{idx+1}:**")
            st.write(msg["content"])

            # Extract fenced code blocks and heuristics (preserved behavior)
            code_blocks = re.findall(r"```(?:python|py)?\s*([\s\S]*?)\s*```", msg["content"], flags=re.IGNORECASE)
            if not code_blocks:
                lines = msg["content"].splitlines()
                code_lines = []; in_code = False
                for line in lines:
                    if re.match(r"\s*(import\s|from\s+\w+\s+import|def\s|class\s|plt\.|sns\.|pd\.|np\.|fig\s*=|px\.)", line):
                        in_code = True
                    if in_code:
                        code_lines.append(line)
                        if line.strip() == "" and len(code_lines) > 3:
                            break
                if code_lines:
                    code_blocks = ["\n".join(code_lines)]

            for i, raw in enumerate(code_blocks):
                code_snippet = raw.strip().replace("```", "").strip()
                if not re.search(r"(?:^|\n)\s*(import\s|from\s+\w+\s+import|def\s|class\s|plt\.|sns\.|pd\.|np\.|print\s*\(|fig\s*=|px\.)", code_snippet):
                    st.caption("Skipped non-code block.")
                    continue
                with st.expander(f"AI Python snippet #{idx+1}-{i+1}", expanded=False):
                    st.code(code_snippet, language="python")
                    if st.button(f"Run snippet #{idx+1}-{i+1}", key=f"run_{idx}_{i}"):
                        out, err = run_generated_code(code_snippet, st.session_state.df, msg_idx=idx)
                        if out:
                            st.subheader("Execution output")
                            st.text(out)
                        if err:
                            st.error("Error running AI code:\n" + str(err))

            # Slide Editor actions: pick sentences or copy (preserved)
            parts = [p.strip() for p in re.split(r'(?<=[\.\?\!])\s+', msg["content"].strip()) if p.strip()]
            if parts:
                chosen = st.multiselect(f"Select sentences/paragraphs from AI #{idx+1} to add to Slide Editor", options=parts, key=f"pick_{idx}")
                c1, c2, c3 = st.columns([1,1,1])
                with c1:
                    if st.button(f"Add selected to Slide Editor (AI #{idx+1})", key=f"add_sel_{idx}"):
                        to_add = "\n\n".join(chosen).strip()
                        if to_add:
                            if ss.slide_editor_text.strip():
                                ss.slide_editor_text = ss.slide_editor_text.rstrip() + "\n\n" + to_add
                            else:
                                ss.slide_editor_text = to_add
                            st.success("Added to Slide Editor.")
                            st.rerun()
                        else:
                            st.warning("No sentences selected.")
                with c2:
                    if st.button(f"Copy full AI #{idx+1} to Slide Editor", key=f"copy_full_{idx}"):
                        if ss.slide_editor_text.strip():
                            ss.slide_editor_text = ss.slide_editor_text.rstrip() + "\n\n" + msg["content"].strip()
                        else:
                            ss.slide_editor_text = msg["content"].strip()
                        st.success("Copied to Slide Editor.")
                        st.rerun()
                with c3:
                    plot_path = ss.msg_plot_map.get(idx)
                    if plot_path:
                        if st.button(f"Use plot from AI #{idx+1} on Slide", key=f"use_plot_{idx}"):
                            ss.slide_editor_selected_plot = plot_path
                            st.success("Plot selected for Slide Editor.")
                    else:
                        st.write("")

    st.markdown("---")
    # Slide Editor (preserved)
    st.subheader("📝 Slide Editor — select or paste text, choose plot")
    ss.slide_editor_title = st.text_input("Slide title", value=ss.slide_editor_title, key="slide_title_input")
    ss.slide_editor_text = st.text_area("Slide content (edit or paste)", value=ss.slide_editor_text, height=220, key="slide_text_area")

    # Show available plots
    st.write("Available plots (from executed code):")
    plots = list(dict.fromkeys(ss.msg_plot_map.values()))  # unique preserve order
    if plots:
        cols = st.columns(min(3, len(plots)))
        for i, p in enumerate(plots):
            with cols[i % max(1, len(cols))]:
                if p and os.path.exists(p):
                    st.image(p, use_column_width=True, caption=f"Plot {i+1}")
                    if st.button(f"Select this plot for slide (#{i+1})", key=f"select_plot_{i}"):
                        ss.slide_editor_selected_plot = p
                        st.success("Plot selected.")
                else:
                    st.write("Missing image")
        if ss.slide_editor_selected_plot:
            if os.path.exists(ss.slide_editor_selected_plot):
                st.info("Current selected plot for slide:")
                st.image(ss.slide_editor_selected_plot, use_column_width=True)
            else:
                ss.slide_editor_selected_plot = None
                st.warning("Previously selected plot missing; cleared.")
    else:
        st.info("No plots available. Run AI code that produces plots to populate this list.")

    ic, cc = st.columns([1,1])
    with ic:
        if st.button("➕ Insert Slide Editor content into PPT", key="insert_slide"):
            title = ss.slide_editor_title or "Slide"
            text = ss.slide_editor_text or ""
            image = ss.slide_editor_selected_plot if ss.slide_editor_selected_plot and os.path.exists(ss.slide_editor_selected_plot) else None
            try:
                add_slide_with_text_and_optional_image(title, text, image_path=image)
                st.success("Slide added to PPT.")
                # optional: clear editor
                ss.slide_editor_text = ""
                ss.slide_editor_selected_plot = None
                ss.slide_editor_title = "Slide Title"
                ss.preview_dirty = True
                st.rerun()
            except Exception as e:
                st.error("Failed to add slide: " + str(e))
    with cc:
        if st.button("🧹 Clear Slide Editor", key="clear_slide_editor"):
            ss.slide_editor_text = ""
            ss.slide_editor_selected_plot = None
            st.success("Slide Editor cleared.")

    st.markdown("---")
    # Save & Download PPT (preserved)
    if st.button("💾 Save & Download PPTX", key="save_download"):
        try:
            tmp = tempfile.NamedTemporaryFile(delete=False, suffix=".pptx")
            ss.ppt.save(tmp.name)
            tmp.close()
            with open(tmp.name, "rb") as f:
                st.download_button("📥 Download PowerPoint", f, file_name="AI_Report_v3.pptx", mime="application/vnd.openxmlformats-officedocument.presentationml.presentation")
        except Exception as e:
            st.error("Failed saving PPTX: " + str(e))

    # follow-up chat input (preserved)
    if st.session_state.df is not None:
        user_input = st.chat_input("Ask a follow-up question...")
        if user_input:
            ss.messages.append({"role": "user", "content": user_input})
            system_instructions = """
            You are an expert process engineer.
            - Always use the user's uploaded dataset (df) provided below.
            - Never create or rely on synthetic/sample datasets.
            - If asked for Python code, reply ONLY with code inside ```python blocks.
            - If asked for conclusions for PPT, provide 1–2 crisp, data-backed sentences and/or a chart.
            """
            if ss.df is not None:
                dataset_context = ss.df.to_csv(index=False)
            else:
                dataset_context = "No dataset available."
            history_text = "\n".join([f"{m['role'].capitalize()}: {m['content']}" for m in ss.messages])
            full_prompt = f"""
            {system_instructions}

            Conversation so far:
            {history_text}

            Dataset (CSV):
            {dataset_context}
            """
            try:
                response = model.generate_content(full_prompt)
                ss.messages.append({"role": "assistant", "content": response.text})
                st.rerun()
            except Exception as e:
                st.error("AI call failed: " + str(e))

# RIGHT column: Live Preview (preserved)
with right:
    st.subheader("🖼 Live Slide Preview")
    c1, c2 = st.columns([1,1])
    with c1:
        if AUTOREFRESH_AVAILABLE:
            ss.autorefresh_on = st.toggle("Auto-refresh", value=ss.autorefresh_on, key="autorefresh_toggle")
        else:
            st.caption("Install streamlit-autorefresh for auto-refresh.")
    with c2:
        ss.refresh_interval_ms = st.number_input("Refresh interval (ms)", min_value=1000, max_value=20000, value=ss.refresh_interval_ms, step=500, key="refresh_interval")

    if AUTOREFRESH_AVAILABLE and ss.autorefresh_on:
        st_autorefresh(interval=ss.refresh_interval_ms, key="live_preview_tick")

    if st.button("🔄 Refresh Preview", key="manual_refresh"):
        ss.preview_dirty = True

    if ss.preview_dirty:
        try:
            ss.preview_images = generate_live_preview_images()
        except Exception as e:
            ss.preview_images = generate_live_preview_images()
        ss.preview_dirty = False

    num_slides = len(ss.preview_images)
    if num_slides == 0:
        ss.slide_idx = 1
        st.info("No slides to preview yet.")
    else:
        ss.slide_idx = min(max(1, ss.slide_idx), num_slides)
        nav1, nav2, nav3 = st.columns([1,6,1])
        with nav1:
            if st.button("◀", key="prev_slide", disabled=(num_slides <= 1 or ss.slide_idx <= 1)):
                ss.slide_idx = max(1, ss.slide_idx - 1)
        with nav2:
            if num_slides >= 2:
                ss.slide_idx = st.slider("Slide", min_value=1, max_value=num_slides, value=ss.slide_idx, key="slide_selector_main")
            else:
                st.markdown("**Slide 1 of 1**")
        with nav3:
            if st.button("▶", key="next_slide", disabled=(num_slides <= 1 or ss.slide_idx >= num_slides)):
                ss.slide_idx = min(num_slides, ss.slide_idx + 1)

        current_img = ss.preview_images[ss.slide_idx - 1]
        st.image(current_img, use_container_width=True, caption=f"Slide {ss.slide_idx} of {num_slides}")

        with st.expander("Thumbnails", expanded=False):
            thumbs = ss.preview_images
            cols = st.columns(3)
            for i, p in enumerate(thumbs):
                with cols[i % 3]:
                    if st.button(f"Go to slide {i+1}", key=f"thumb_{i}"):
                        ss.slide_idx = i+1
                    st.image(p, use_column_width=True, caption=f"{i+1}")

# Floating viewer pinned
_render_floating_viewer(ss.preview_images, ss.slide_idx)
